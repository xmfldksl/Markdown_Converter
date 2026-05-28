import os # 파일 경로 처리를 위한 내장 모듈을 가져옵니다.
from pptx import Presentation # 파워포인트 파일을 읽고 해독하기 위한 외부 라이브러리 클래스를 가져옵니다.

def extract_pptx_content(file_path, progress_callback=None):
    # 파워포인트 파일 경로를 입력받아 슬라이드 내부 텍스트와 표를 마크다운으로 변환하는 함수입니다.
    md_output = "" # 최종 마크다운 결과물이 누적될 빈 텍스트 문자열을 초기화합니다.
    
    prs = Presentation(file_path) # 전달받은 경로의 파워포인트 파일을 메모리에 객체로 불러옵니다.
    total_slides = len(prs.slides) # 파일에 포함된 전체 슬라이드 개수를 파악하여 변수에 저장합니다.
    
    for i, slide in enumerate(prs.slides): # 첫 번째 슬라이드부터 마지막 슬라이드까지 순서대로 반복 탐색합니다.
        md_output += f"## Slide {i + 1}\n\n" # 현재 슬라이드 번호를 마크다운 중제목 형태로 문자열에 추가합니다.
        
        for shape in slide.shapes: # 현재 슬라이드 내부에 배치된 모든 도형, 텍스트 상자, 표 객체들을 순회합니다.
            if shape.has_text_frame: # 만약 현재 객체가 글자를 포함할 수 있는 텍스트 프레임(텍스트 상자 등)이라면.
                text = shape.text.strip() # 프레임 내부의 전체 텍스트를 추출하고 양끝의 불필요한 공백을 제거합니다.
                if text: # 텍스트가 텅 비어있지 않고 실제 내용이 존재한다면.
                    md_output += text + "\n\n" # 추출된 텍스트를 결과물에 이어 붙이고 단락 구분을 위한 줄바꿈을 넣습니다.
                    
            elif shape.has_table: # 만약 현재 객체가 글자 상자가 아니라 데이터가 들어있는 표(테이블) 객체라면.
                table = shape.table # 표 객체 데이터를 다루기 쉽게 변수에 할당합니다.
                for row_idx, row in enumerate(table.rows): # 표의 첫 번째 가로줄(행)부터 순서대로 반복 탐색합니다.
                    # 행 내부의 각 칸(셀)에 접근하여 텍스트를 추출하고, 마크다운 표 문법이 깨지지 않도록 줄바꿈과 파이프 기호를 치환하여 리스트로 묶습니다.
                    row_data = [cell.text_frame.text.replace('\n', '<br>').replace('|', '\\|').strip() for cell in row.cells]
                    md_output += "| " + " | ".join(row_data) + " |\n" # 치환된 칸 데이터들을 파이프 기호로 연결하여 마크다운 표의 한 줄을 완성합니다.
                    
                    if row_idx == 0: # 방금 처리한 줄이 표의 가장 최상단 제목 행이었다면.
                        md_output += "| " + " | ".join(["---"] * len(row.cells)) + " |\n" # 칸 개수만큼 대시 기호를 넣어 마크다운 표의 제목 구분선을 삽입합니다.
                        
                md_output += "\n" # 하나의 표 처리가 모두 끝났으므로 다른 요소와 구분하기 위해 줄바꿈을 하나 추가합니다.
        
        md_output += "---\n\n" # 한 슬라이드의 모든 내용 추출이 완료되었음을 나타내는 마크다운 수평 구분선을 넣습니다.
        
        if progress_callback: # UI 화면으로 진행 상황을 보고할 콜백 함수가 인자로 전달되었다면.
            percent = int(((i + 1) / total_slides) * 100) # 현재까지 처리한 슬라이드 비율을 정수형 백분율로 계산합니다.
            progress_callback(percent) # 계산된 퍼센트 정수 숫자를 메인 스레드로 전송하여 화면을 갱신합니다.
            
    return md_output # 모든 슬라이드 분석이 끝나고 완성된 전체 마크다운 문자열을 호출자에게 반환합니다.